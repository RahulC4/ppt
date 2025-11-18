# ============================================================
# generate_ppt.py  (FINAL – FIXED & STABLE)
# ============================================================

import os
import tempfile
import uuid
import json
import re
import requests
from pptx import Presentation
from pptx.util import Inches, Pt
from utils import (
    get_env, safe_json_load, logger, now_ts,
    ensure_dir, text_client, image_client
)
from search_utils import semantic_search
from azure_blob_utils import upload_ppt_to_blob, upload_json_to_blob


# Ensure directory exists for design JSONs
ensure_dir("design_jsons")

CHAT_MODEL = get_env("CHAT_MODEL", required=True)
IMAGE_MODEL = get_env("IMAGE_MODEL", required=True)


# ---------------------------------------------------------
# AUTO–DETECTION: slide count + theme
# ---------------------------------------------------------
def parse_user_intent(prompt: str):
    num_slides = None
    theme = None

    m = re.search(r"(\d+)\s+slides?", prompt.lower())
    if m:
        num_slides = int(m.group(1))

    theme_words = [
        "corporate", "modern", "minimal", "professional",
        "dark", "light", "colorful", "gradient", "flat"
    ]
    for t in theme_words:
        if t in prompt.lower():
            theme = t.capitalize()
            break

    return num_slides, theme


# ---------------------------------------------------------
# GPT PLAN GENERATOR
# ---------------------------------------------------------
def call_llm_plan(prompt, style, design_context, references_text,
                  num_slides=None, theme=None):

    # Safety if empty
    if not design_context:
        design_context = ["no_design_data"]
    if not references_text:
        references_text = ["no_reference_text"]

    sys_prompt = (
        "You are an expert presentation designer.\n"
        "Create a structured slide plan.\n"
        "Return ONLY valid JSON as a list of slides:\n"
        "[{\"title\": str, \"bullets\": [str], \"visual_required\": bool, \"visual_prompt\": str}]\n\n"
        "If the user explicitly asks for images, ALWAYS set visual_required to true.\n"
        "Respond ONLY with JSON. No explanations.\n"
    )

    if theme:
        sys_prompt += f"\nTheme preference: {theme}\n"

    # Truncate just to keep prompt stable
    sys_prompt += f"\nDesign Context: {str(design_context)[:1500]}"
    sys_prompt += f"\nReferences: {str(references_text)[:1500]}\n"

    user_prompt = f"Generate a presentation plan for: {prompt} (Style: {style})"
    if num_slides:
        user_prompt += f". Use exactly {num_slides} slides."

    try:
        resp = text_client.chat.completions.create(
            model=CHAT_MODEL,
            messages=[
                {"role": "system", "content": sys_prompt},
                {"role": "user", "content": user_prompt}
            ],
            max_completion_tokens=1200,
            temperature=0.7
        )
        raw = resp.choices[0].message.content
        plan = safe_json_load(raw)

        if not plan:
            logger.warning("Invalid JSON from GPT. Retrying...")
            resp2 = text_client.chat.completions.create(
                model=CHAT_MODEL,
                messages=[
                    {"role": "system", "content": sys_prompt},
                    {"role": "user", "content": user_prompt}
                ],
                max_completion_tokens=1200,
                temperature=0.7
            )
            plan = safe_json_load(resp2.choices[0].message.content)

        if not plan:
            return [{"title": "Intro", "bullets": ["Overview"], "visual_required": True}]

        return plan

    except Exception as e:
        logger.exception(f"Plan generation failed: {e}")
        return [{"title": "Intro", "bullets": ["Overview"], "visual_required": True}]


# ---------------------------------------------------------
# VISUAL GENERATOR (Image model)
# ---------------------------------------------------------
def generate_visual_image(prompt: str):
    # Keep prompt short for Azure image model
    prompt = prompt[:250]

    try:
        resp = image_client.images.generate(
            model=IMAGE_MODEL,
            prompt=prompt,
            size="1024x1024"
        )

        url = resp.data[0].url
        img_bytes = requests.get(url, timeout=20).content

        tmp = tempfile.NamedTemporaryFile(delete=False, suffix=".png")
        tmp.write(img_bytes)
        tmp.close()

        return tmp.name

    except Exception as e:
        logger.warning(f"Image generation failed for '{prompt}': {e}")
        return None


# ---------------------------------------------------------
# PPT BUILDER
# ---------------------------------------------------------
def build_ppt(slides):
    prs = Presentation()

    for s in slides:
        layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(layout)

        # Title
        try:
            slide.shapes.title.text = s.get("title", "")
        except:
            pass

        # Bullets
        try:
            tf = slide.placeholders[1].text_frame
            tf.clear()
            for b in s.get("bullets", []):
                p = tf.add_paragraph()
                p.text = b
                p.font.size = Pt(18)
        except:
            pass

        # Image
        if s.get("image_path"):
            try:
                slide.shapes.add_picture(
                    s["image_path"],
                    Inches(0.5), Inches(3.0),
                    width=Inches(8)
                )
            except Exception as e:
                logger.debug(f"Failed to add image: {e}")

    out = os.path.join(
        tempfile.gettempdir(),
        f"generated_presentation_{uuid.uuid4().hex[:8]}.pptx"
    )
    prs.save(out)
    return out


# ---------------------------------------------------------
# MAIN PIPELINE
# ---------------------------------------------------------
def generate_presentation(prompt: str, style="Auto", requested_num_slides=None,
                          theme=None, tag_filters=None):

    detected_slides, detected_theme = parse_user_intent(prompt)
    requested_num_slides = requested_num_slides or detected_slides
    theme = theme or detected_theme

    # Reduce top_k = 5
    refs = semantic_search(prompt, top_k=5, tags=tag_filters) or []

    design_context = []
    reference_text = []

    for r in refs:
        ppt_name = r.get("ppt_name")
        snippet = (r.get("text") or "")[:400]
        if snippet:
            reference_text.append(snippet)

        json_path = os.path.join("design_jsons", os.path.basename(ppt_name) + ".json")
        if os.path.exists(json_path):
            try:
                with open(json_path, "r", encoding="utf-8") as f:
                    design_context.append(json.load(f))
            except:
                pass

    logger.info(
        f"Loaded {len(design_context)} design JSONs and {len(reference_text)} text snippets."
    )

    # Build plan
    plan = call_llm_plan(
        prompt,
        style,
        design_context,
        reference_text,
        num_slides=requested_num_slides,
        theme=theme
    )

    # Force images if user asked
    force_images = "image" in prompt.lower() or "images" in prompt.lower()

    # Build slide objects
    slides = []
    for sp in plan:
        slide = {
            "title": sp.get("title", "Untitled"),
            "bullets": sp.get("bullets", []),
            "image_path": None
        }

        # Always use a SHORT image prompt
        img_prompt = f"Professional illustration representing {slide['title']}."
        img_prompt = img_prompt[:200]

        if sp.get("visual_required", False) or force_images:
            slide["image_path"] = generate_visual_image(img_prompt)

        slides.append(slide)

    # Build PPT
    out_path = build_ppt(slides)

    # Upload to Azure Blob
    fname = f"generated_{uuid.uuid4().hex[:8]}.pptx"
    upload_ppt_to_blob(out_path, fname)

    # Log metadata
    log = {
        "timestamp": now_ts(),
        "prompt": prompt,
        "style": style,
        "theme": theme,
        "requested_num_slides": requested_num_slides,
        "references_used": len(reference_text),
        "design_jsons_used": len(design_context),
        "slides_generated": len(slides),
        "ppt_file": fname
    }

    upload_json_to_blob(
        json.dumps(log, indent=2).encode("utf-8"),
        f"logs/{fname}.json"
    )

    return out_path, log
