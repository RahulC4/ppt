# generate_ppt.py – FINAL CLEAN VERSION
import os
import tempfile
import uuid
import json
import re
import requests
from pptx import Presentation
from pptx.util import Inches, Pt

from utils import (
    text_client,
    image_client,
    embedding_client,
    CHAT_MODEL,
    EMBEDDING_MODEL,
    IMAGE_MODEL,
    safe_json_load,
    logger,
    now_ts,
    ensure_dir
)

from search_utils import semantic_search
from azure_blob_utils import upload_ppt_to_blob, upload_json_to_blob


# ensure design json directory
ensure_dir("design_jsons")


# ---------------------------------------------------------
# Extract slide count + theme from natural prompt
# ---------------------------------------------------------
def parse_user_intent(prompt: str):
    num_slides = None
    theme = None

    match = re.search(r"(\d+)\s+slides?", prompt.lower())
    if match:
        num_slides = int(match.group(1))

    theme_keywords = [
        "modern", "minimal", "corporate", "professional",
        "dark", "light", "colorful", "gradient"
    ]
    for w in theme_keywords:
        if w in prompt.lower():
            theme = w.capitalize()
            break

    return num_slides, theme


# ---------------------------------------------------------
# LLM: create plan + visuals decisions
# ---------------------------------------------------------
def call_llm_plan(prompt, style, design_context, references_text, num_slides=None, theme=None):
    sys_prompt = (
        "You are an expert presentation designer. "
        "Output ONLY a JSON list of slides.\n\n"
        "Each slide must follow this format:\n"
        "{"
        "\"title\": str, "
        "\"bullets\": [str], "
        "\"visual_required\": bool, "
        "\"visual_prompt\": str"
        "}\n\n"
        "If reference texts AND design JSON match, produce visuals. "
        "If not, decide on visuals independently.\n"
    )

    if theme:
        sys_prompt += f"\nPreferred theme: {theme}.\n"

    sys_prompt += f"\nDesign Context (truncated): {json.dumps(design_context)[:3500]}"
    sys_prompt += f"\nReference snippets: {json.dumps(references_text)[:2000]}\n"

    user_prompt = f"Create a PowerPoint plan for: {prompt}. Style: {style}."
    if num_slides:
        user_prompt += f" Produce exactly {num_slides} slides."

    try:
        resp = text_client.chat.completions.create(
            model=CHAT_MODEL,
            messages=[
                {"role": "system", "content": sys_prompt},
                {"role": "user", "content": user_prompt}
            ]
        )
        text = resp.choices[0].message.content
        plan = safe_json_load(text)

        if not plan:
            logger.warning("LLM returned invalid JSON. Using fallback slide.")
            return [
                {
                    "title": "Introduction",
                    "bullets": ["Overview"],
                    "visual_required": False
                }
            ]

        return plan

    except Exception as e:
        logger.exception(f"Plan generation failed: {e}")
        return [
            {
                "title": "Introduction",
                "bullets": ["Overview"],
                "visual_required": False
            }
        ]


# ---------------------------------------------------------
# Visual generator using Image client
# ---------------------------------------------------------
def generate_visual_image(prompt):
    try:
        resp = image_client.images.generate(
            model=IMAGE_MODEL,
            prompt=prompt,
            size="1024x1024"
        )
        url = resp.data[0].url
        r = requests.get(url, timeout=20)

        tmp = tempfile.NamedTemporaryFile(delete=False, suffix=".png")
        tmp.write(r.content)
        tmp.close()
        return tmp.name

    except Exception as e:
        logger.warning(f"Visual generation failed: {e}")
        return None


# ---------------------------------------------------------
# Build PPTX
# ---------------------------------------------------------
def build_ppt(slides):
    prs = Presentation()
    for s in slides:
        try:
            layout = prs.slide_layouts[1]
        except:
            layout = prs.slide_layouts[0]

        slide = prs.slides.add_slide(layout)

        # title
        try:
            slide.shapes.title.text = s["title"]
        except:
            pass

        # bullets
        try:
            tf = slide.placeholders[1].text_frame
            tf.clear()
            for b in s["bullets"]:
                p = tf.add_paragraph()
                p.text = b
                p.level = 0
                p.font.size = Pt(18)
        except:
            pass

        # image (if any)
        if s.get("image_path"):
            try:
                slide.shapes.add_picture(
                    s["image_path"],
                    Inches(0.5), Inches(3),
                    width=Inches(8)
                )
            except Exception as e:
                logger.debug(f"Failed to add image: {e}")

    out_path = os.path.join(
        tempfile.gettempdir(),
        f"generated_{uuid.uuid4().hex[:8]}.pptx"
    )
    prs.save(out_path)
    return out_path


# ---------------------------------------------------------
# MAIN PIPELINE
# ---------------------------------------------------------
def generate_presentation(
    prompt: str,
    style: str = "Auto",
    requested_num_slides: int = None,
    theme: str = None,
    tag_filters: list = None,
    tags: list = None   # legacy param support
):
    # backwards compatible arg
    if tag_filters is None and tags is not None:
        tag_filters = tags

    # detect in prompt
    detected_slides, detected_theme = parse_user_intent(prompt)
    if not requested_num_slides:
        requested_num_slides = detected_slides
    if not theme:
        theme = detected_theme

    # retrieve from chroma
    refs = semantic_search(prompt, top_k=10, tags=tag_filters) or []

    design_context = []
    reference_texts = []

    for r in refs:
        ppt_name = r.get("ppt_name")
        snippet = (r.get("text") or "")[:400]
        if snippet:
            reference_texts.append(snippet)

        json_path = os.path.join("design_jsons", f"{ppt_name}.json")
        if os.path.exists(json_path):
            try:
                with open(json_path, "r", encoding="utf-8") as f:
                    design_context.append(json.load(f))
            except:
                pass

    logger.info(f"Loaded {len(design_context)} design JSONs and {len(reference_texts)} reference snippets.")

    # ask LLM for plan
    plan = call_llm_plan(
        prompt,
        style,
        design_context,
        reference_texts,
        num_slides=requested_num_slides,
        theme=theme
    )

    # generate visuals & combine
    slides = []
    for p in plan:
        slide = {
            "title": p.get("title", "Untitled"),
            "bullets": p.get("bullets", []),
            "image_path": None
        }

        if p.get("visual_required"):
            img = generate_visual_image(p.get("visual_prompt", slide["title"]))
            if img:
                slide["image_path"] = img

        slides.append(slide)

    # build PPT
    out_path = build_ppt(slides)

    # upload
    fname = f"generated_{uuid.uuid4().hex[:8]}.pptx"
    upload_ppt_to_blob(out_path, fname)

    # upload logs
    log = {
        "timestamp": now_ts(),
        "prompt": prompt,
        "slides_generated": len(slides),
        "references_used": len(reference_texts),
        "design_jsons_used": len(design_context),
        "file": fname
    }
    upload_json_to_blob(
        json.dumps(log, indent=2).encode("utf-8"),
        f"logs/{fname}.json"
    )

    return out_path, log
