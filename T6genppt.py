# ============================================================
# generate_ppt.py  (DEBUG-FRIENDLY fixed version)
# ============================================================

import os
import tempfile
import uuid
import json
import re
import base64
from pptx import Presentation
from pptx.util import Inches, Pt
from PIL import Image
from utils import (
    get_env, safe_json_load, logger, now_ts,
    ensure_dir, text_client, image_client
)
from search_utils import semantic_search
from azure_blob_utils import upload_ppt_to_blob, upload_json_to_blob

ensure_dir("design_jsons")

CHAT_MODEL = get_env("CHAT_MODEL", required=True)
IMAGE_MODEL = get_env("IMAGE_MODEL", required=True)


# ---------- helpers ----------
def parse_user_intent(prompt: str):
    num_slides = None
    theme = None
    m = re.search(r"(\d+)\s+slides?", prompt.lower())
    if m:
        num_slides = int(m.group(1))
    for t in ["corporate", "modern", "minimal", "professional", "dark", "light"]:
        if t in prompt.lower():
            theme = t.capitalize()
            break
    return num_slides, theme


def call_llm_plan(prompt, style, design_context, references_text,
                  num_slides=None, theme=None):
    sys_prompt = (
        "You are an expert management consulting presentation designer.\n"
        "Return STRICT JSON ONLY:\n"
        "[ {\"title\": str, \"bullets\": [str], \"visual_required\": bool, \"visual_prompt\": str } ]\n"
        "If the user asks for images, mark visual_required=true.\n"
        "Do NOT include any text inside image prompts.\n"
    )
    if theme:
        sys_prompt += f"\nTheme: {theme}\n"

    sys_prompt += f"\nDesign Context: {json.dumps(design_context)[:3000]}"
    sys_prompt += f"\nReference Snippets: {json.dumps(references_text)[:1500]}\n"

    user_prompt = f"Create a {style} PowerPoint plan for: {prompt}"
    if num_slides:
        user_prompt += f". Create exactly {num_slides} slides."

    try:
        resp = text_client.chat.completions.create(
            model=CHAT_MODEL,
            messages=[{"role": "system", "content": sys_prompt},
                      {"role": "user", "content": user_prompt}],
            max_completion_tokens=1500,
            temperature=1
        )
        raw = resp.choices[0].message.content
        plan = safe_json_load(raw)
        if not plan:
            logger.warning("Invalid plan JSON received from LLM; retrying once.")
            resp2 = text_client.chat.completions.create(
                model=CHAT_MODEL,
                messages=[{"role": "system", "content": sys_prompt},
                          {"role": "user", "content": user_prompt}],
                max_completion_tokens=1200,
                temperature=0.7
            )
            plan = safe_json_load(resp2.choices[0].message.content)

        if not plan:
            logger.warning("LLM did not return JSON plan. Using fallback single intro slide.")
            return [{"title": "Intro", "bullets": ["Overview"], "visual_required": False}]

        return plan

    except Exception as e:
        logger.exception("Plan generation failed")
        return [{"title": "Intro", "bullets": ["Overview"], "visual_required": False}]


def generate_visual_image(prompt: str):
    # Use provided prompt if present, but force no-text clause
    img_prompt = (prompt or "") + " Minimal, clean, professional illustration. No text, no labels."
    # ensure we don't create unbounded prompts (but we do not cut off important info)
    try:
        resp = image_client.images.generate(
            model=IMAGE_MODEL,
            prompt=img_prompt,
            size="1024x1024"
        )
        # handle base64 output
        image_b64 = getattr(resp.data[0], "b64_json", None)
        if image_b64 is None:
            # older behavior (url) or different schema -> try to pull url
            url = getattr(resp.data[0], "url", None)
            if url:
                import requests
                tmp = tempfile.NamedTemporaryFile(delete=False, suffix=".png")
                tmp.write(requests.get(url, timeout=20).content)
                tmp.close()
                return tmp.name
            logger.warning("Image response contained no base64 nor url.")
            return None

        img_bytes = base64.b64decode(image_b64)
        tmp = tempfile.NamedTemporaryFile(delete=False, suffix=".png")
        tmp.write(img_bytes)
        tmp.close()
        return tmp.name

    except Exception as e:
        logger.exception("Image generation failed")
        return None


def build_ppt(slides):
    prs = Presentation()
    for s in slides:
        try:
            layout = prs.slide_layouts[1]
            slide = prs.slides.add_slide(layout)

            # title
            try:
                slide.shapes.title.text = s.get("title", "")
            except Exception:
                pass

            # body text
            try:
                body = slide.placeholders[1]
                tf = body.text_frame
                tf.clear()
                for b in s.get("bullets", []):
                    p = tf.add_paragraph()
                    p.text = b
                    p.font.size = Pt(18)
            except Exception:
                pass

            # If no image -> keep full width but nudge below title
            if not s.get("image_path"):
                try:
                    body = slide.placeholders[1]
                    body.top = body.top + Inches(0.2)
                except Exception:
                    pass
                continue

            # Shrink body to leave right column only when image exists
            try:
                body = slide.placeholders[1]
                body.width = prs.slide_width - Inches(4.0)
                body.top = body.top + Inches(0.3)
            except Exception:
                pass

            # Insert image right column
            try:
                img_path = s["image_path"]
                img = Image.open(img_path)
                iw, ih = img.size
                aspect = iw / ih if ih != 0 else 1.0

                max_w = Inches(3.0)
                max_h = Inches(2.5)
                if aspect >= 1:
                    final_w = max_w
                    final_h = final_w / aspect
                    if final_h > max_h:
                        final_h = max_h
                        final_w = final_h * aspect
                else:
                    final_h = max_h
                    final_w = final_h * aspect
                    if final_w > max_w:
                        final_w = max_w
                        final_h = final_w / aspect

                slide_w = prs.slide_width
                left = slide_w - final_w - Inches(0.5)
                title_shape = slide.shapes.title
                top = title_shape.top + title_shape.height + Inches(0.2)
                slide.shapes.add_picture(img_path, left, top, width=final_w, height=final_h)
            except Exception:
                logger.exception("Failed placing image on slide")

        except Exception:
            logger.exception("Error while building slide (continuing with next slide)")
            continue

    out = os.path.join(tempfile.gettempdir(), f"generated_presentation_{uuid.uuid4().hex[:8]}.pptx")
    prs.save(out)
    return out


def generate_presentation(prompt: str, style="Auto", requested_num_slides=None,
                          theme=None, tag_filters=None):
    logger.info("generate_presentation called")
    try:
        detected_slides, detected_theme = parse_user_intent(prompt)
        requested_num_slides = requested_num_slides or detected_slides
        theme = theme or detected_theme

        refs = semantic_search(prompt, top_k=5, tags=tag_filters) or []

        design_context = []
        reference_text = []

        for r in refs:
            ppt_name = r.get("ppt_name")
            txt = (r.get("text") or "")[:500]
            if txt:
                reference_text.append(txt)
            json_path = os.path.join("design_jsons", os.path.basename(ppt_name) + ".json")
            if os.path.exists(json_path):
                try:
                    with open(json_path, "r", encoding="utf-8") as f:
                        design_context.append(json.load(f))
                except Exception:
                    logger.warning(f"Failed to load design json {json_path}")

        logger.info(f"Loaded {len(design_context)} design JSONs and {len(reference_text)} text snippets.")

        # ask LLM for plan
        plan = call_llm_plan(prompt, style, design_context, reference_text,
                             num_slides=requested_num_slides, theme=theme)

        # force images if user explicitly asked
        force_images = "image" in prompt.lower() or "images" in prompt.lower()

        slides = []
        for sp in plan:
            slide = {
                "title": sp.get("title", "Untitled"),
                "bullets": sp.get("bullets", []),
                "image_path": None
            }
            if sp.get("visual_required") or force_images:
                try:
                    slide["image_path"] = generate_visual_image(sp.get("visual_prompt", ""))
                except Exception:
                    logger.exception("generate_visual_image failed for slide")
                    slide["image_path"] = None
            slides.append(slide)

        out_path = build_ppt(slides)

        fname = f"generated_{uuid.uuid4().hex[:8]}.pptx"
        upload_ppt_to_blob(out_path, fname)

        log = {
            "timestamp": now_ts(),
            "prompt": prompt,
            "slides_generated": len(slides),
            "ppt_file": fname
        }
        upload_json_to_blob(json.dumps(log, indent=2).encode("utf-8"), f"logs/{fname}.json")
        return out_path, log

    except Exception:
        logger.exception("generate_presentation failed")
        raise
