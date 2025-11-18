# ============================================================
# generate_ppt.py  (FINAL – FULLY WORKING WITH BASE64 IMAGES)
# ============================================================

import os
import tempfile
import uuid
import json
import re
import requests
import base64
from pptx import Presentation
from pptx.util import Inches, Pt
from utils import (
    get_env, safe_json_load, logger, now_ts,
    ensure_dir, text_client, image_client
)
from search_utils import semantic_search
from azure_blob_utils import upload_ppt_to_blob, upload_json_to_blob


# Ensure design folder exists
ensure_dir("design_jsons")

CHAT_MODEL = get_env("CHAT_MODEL", required=True)
IMAGE_MODEL = get_env("IMAGE_MODEL", required=True)


# ---------------------------------------------------------
# AUTO DETECT USER INTENT
# ---------------------------------------------------------
def parse_user_intent(prompt: str):
    num_slides = None
    theme = None

    m = re.search(r"(\d+)\s+slides?", prompt.lower())
    if m:
        num_slides = int(m.group(1))

    theme_words = [
        "corporate", "modern", "minimal", "professional",
        "dark", "light", "colorful", "gradient", "flat"
    ]
    for t in theme_words:
        if t in prompt.lower():
            theme = t.capitalize()
            break

    return num_slides, theme


# ---------------------------------------------------------
# GPT PLAN GENERATOR (Safe JSON)
# ---------------------------------------------------------
def call_llm_plan(prompt, style, design_context, references_text,
                  num_slides=None, theme=None):

    if not design_context:
        design_context = ["no_design_data"]
    if not references_text:
        references_text = ["no_reference_text"]

    sys_prompt = (
        "You are an expert presentation designer.\n"
        "Create a structured slide plan.\n"
        "Return ONLY valid JSON in this list format:\n"
        "[{\"title\": str, \"bullets\": [str], \"visual_required\": bool, \"visual_prompt\": str}]\n"
        "Use the design JSONs for style cues and reference snippets for content.\n"
        "If both match well, enable visuals. Otherwise, decide independently.\n"
        "If slide count or theme is detected from user text, follow it.\n"
        "If the user explicitly asks for images, ALWAYS set visual_required to true.\n"
        "Do NOT include explanations. Output JSON only.\n"
    )

    if theme:
        sys_prompt += f"\nTheme preference: {theme}"

    sys_prompt += f"\nDesign Context: {str(design_context)[:1200]}"
    sys_prompt += f"\nReference Snippets: {str(references_text)[:1200]}\n"

    user_prompt = f"Generate a presentation plan for: {prompt} (Style={style})"
    if num_slides:
        user_prompt += f". Use exactly {num_slides} slides."

    try:
        resp = text_client.chat.completions.create(
            model=CHAT_MODEL,
            messages=[
                {"role": "system", "content": sys_prompt},
                {"role": "user", "content": user_prompt}
            ],
            max_completion_tokens=1200,
            temperature=0.7
        )
        raw = resp.choices[0].message.content
        plan = safe_json_load(raw)

        if not plan:
            logger.warning("Invalid plan JSON, retrying...")
            resp2 = text_client.chat.completions.create(
                model=CHAT_MODEL,
                messages=[
                    {"role": "system", "content": sys_prompt},
                    {"role": "user", "content": user_prompt}
                ],
                max_completion_tokens=1200,
                temperature=0.7
            )
            plan = safe_json_load(resp2.choices[0].message.content)

        if not plan:
            return [{"title": "Intro", "bullets": ["Overview"], "visual_required": True}]

        return plan

    except Exception as e:
        logger.exception(f"Plan generation failed: {e}")
        return [{"title": "Intro", "bullets": ["Overview"], "visual_required": True}]


# ---------------------------------------------------------
# IMAGE GENERATOR (Handles Azure Base64 Response)
# ---------------------------------------------------------
def generate_visual_image(prompt: str):
    prompt = prompt[:250]  # Keep prompt short

    try:
        resp = image_client.images.generate(
            model=IMAGE_MODEL,
            prompt=prompt,
            size="1024x1024"
        )

        # Azure returns base64 JSON (b64_json field)
        image_b64 = resp.data[0].b64_json
        img_bytes = base64.b64decode(image_b64)

        tmp = tempfile.NamedTemporaryFile(delete=False, suffix=".png")
        tmp.write(img_bytes)
        tmp.close()
        return tmp.name

    except Exception as e:
        logger.warning(f"Image generation failed for '{prompt}': {e}")
        return None


# ---------------------------------------------------------
# PPT BUILDER
# ---------------------------------------------------------
def build_ppt(slides):
    prs = Presentation()

    for s in slides:
        layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(layout)

        try:
            slide.shapes.title.text = s.get("title", "")
        except:
            pass

        try:
            tf = slide.placeholders[1].text_frame
            tf.clear()
            for b in s.get("bullets", []):
                p = tf.add_paragraph()
                p.text = b
                p.font.size = Pt(18)
        except:
            pass

        if s.get("image_path"):
            try:
                slide.shapes.add_picture(
                    s["image_path"],
                    Inches(0.5),
                    Inches(3.0),
                    width=Inches(8)
                )
            except Exception as e:
                logger.debug(f"Failed to add image: {e}")

    out = os.path.join(
        tempfile.gettempdir(),
        f"generated_presentation_{uuid.uuid4().hex[:8]}.pptx"
    )
    prs.save(out)
    return out


# ---------------------------------------------------------
# MAIN PIPELINE
# ---------------------------------------------------------
def generate_presentation(prompt: str, style="Auto", requested_num_slides=None,
                          theme=None, tag_filters=None):

    detected_slides, detected_theme = parse_user_intent(prompt)
    requested_num_slides = requested_num_slides or detected_slides
    theme = theme or detected_theme

    refs = semantic_search(prompt, top_k=5, tags=tag_filters) or []

    design_context = []
    reference_text = []

    for r in refs:
        snippet = (r.get("text") or "")[:300]
        if snippet:
            reference_text.append(snippet)

        ppt_name = r.get("ppt_name")
        json_path = os.path.join("design_jsons", os.path.basename(ppt_name) + ".json")
        if os.path.exists(json_path):
            try:
                with open(json_path, "r", encoding="utf-8") as f:
                    design_context.append(json.load(f))
            except:
                pass

    logger.info(f"Loaded {len(design_context)} design JSONs and {len(reference_text)} text snippets.")

    plan = call_llm_plan(
        prompt,
        style,
        design_context,
        reference_text,
        num_slides=requested_num_slides,
        theme=theme
    )

    force_images = "image" in prompt.lower() or "images" in prompt.lower()

    slides = []
    for sp in plan:
        slide = {
            "title": sp.get("title", "Untitled"),
            "bullets": sp.get("bullets", []),
            "image_path": None
        }

        # Short, stable prompt
        img_prompt = f"Professional illustration representing {slide['title']}."

        if sp.get("visual_required", False) or force_images:
            slide["image_path"] = generate_visual_image(img_prompt)

        slides.append(slide)

    out_path = build_ppt(slides)

    fname = f"generated_{uuid.uuid4().hex[:8]}.pptx"
    upload_ppt_to_blob(out_path, fname)

    log = {
        "timestamp": now_ts(),
        "prompt": prompt,
        "style": style,
        "theme": theme,
        "requested_num_slides": requested_num_slides,
        "references_used": len(reference_text),
        "design_jsons_used": len(design_context),
        "slides_generated": len(slides),
        "ppt_file": fname
    }

    upload_json_to_blob(
        json.dumps(log, indent=2).encode("utf-8"),
        f"logs/{fname}.json"
    )

    return out_path, log
