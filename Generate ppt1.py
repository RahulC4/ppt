# generate_ppt.py
import os
import tempfile
import uuid
import json
import re
import requests
from pptx import Presentation
from pptx.util import Inches, Pt
from openai import AzureOpenAI
from utils import get_env, safe_json_load, logger, now_ts, ensure_dir
from search_utils import semantic_search
from azure_blob_utils import upload_ppt_to_blob, upload_json_to_blob

# ensure design json dir exists
try:
    ensure_dir("design_jsons")
except Exception:
    os.makedirs("design_jsons", exist_ok=True)

# === Clients: text_client (chat + embeddings) and image_client (images) ===
text_client = AzureOpenAI(
    azure_endpoint=get_env("OPENAI_API_BASE", required=True),
    api_key=get_env("OPENAI_API_KEY", required=True),
    api_version=get_env("OPENAI_API_VERSION", "2024-05-01-preview")
)

image_client = AzureOpenAI(
    azure_endpoint=get_env("IMAGE_API_BASE", default=get_env("OPENAI_API_BASE")),
    api_key=get_env("IMAGE_API_KEY", default=get_env("OPENAI_API_KEY")),
    api_version=get_env("OPENAI_API_VERSION", "2024-05-01-preview")
)

CHAT_MODEL = get_env("CHAT_MODEL", "gpt-5-mini")
EMBEDDING_MODEL = get_env("EMBEDDING_MODEL", "text-embedding-3-large")
DALLE_MODEL = get_env("DALLE_MODEL", "gpt-image-1-mini")

# ---------------------------------------------------------
# Helper: detect slide count and theme from natural prompt
# ---------------------------------------------------------
def parse_user_intent(prompt: str):
    """Extract number of slides and possible theme hints from user prompt."""
    num_slides = None
    theme = None

    # numeric detection like "5 slides"
    m = re.search(r'(\d+)\s+slides?', prompt.lower())
    if m:
        try:
            num_slides = int(m.group(1))
        except Exception:
            num_slides = None

    # word-based small numbers (one..ten) - optional
    words_to_nums = {
        "one":1,"two":2,"three":3,"four":4,"five":5,
        "six":6,"seven":7,"eight":8,"nine":9,"ten":10
    }
    if num_slides is None:
        for w, n in words_to_nums.items():
            if re.search(rf'\b{w}\b.*slides?', prompt.lower()):
                num_slides = n
                break

    # theme keywords
    theme_keywords = ["modern", "minimal", "corporate", "professional", "dark", "light", "colorful", "flat", "gradient"]
    for word in theme_keywords:
        if word in prompt.lower():
            theme = word.capitalize()
            break

    return num_slides, theme

# ---------------------------------------------------------
# Helper: call LLM to create slide plan (with visual decisions)
# ---------------------------------------------------------
def call_llm_plan(prompt: str, style: str, design_context: list, references_text: list,
                  num_slides: int = None, theme: str = None):
    """
    Ask LLM to produce a slide plan. Response MUST be JSON array of slides:
    [{"title": "...", "bullets": ["..."], "visual_required": true/false, "visual_prompt": "..."}]
    """
    sys_prompt = (
        "You are an expert presentation designer. Output ONLY valid JSON: a list of slides. "
        "Each slide must be an object with keys: "
        "{\"title\": str, \"bullets\": [str], \"visual_required\": bool, \"visual_prompt\": str (optional)}.\n\n"
        "Use the provided reference snippets for content and the provided design JSONs for style/visual cues. "
        "If references align in both content and design, include visuals; otherwise decide visuals yourself. "
        "If a specific number of slides is requested, produce exactly that number. Be concise."
    )

    if theme:
        sys_prompt += f"\nTheme preference: {theme}."

    # trim contexts for token safety
    try:
        design_excerpt = json.dumps(design_context, indent=None)[:4000]
    except Exception:
        design_excerpt = "[]"
    try:
        refs_excerpt = json.dumps(references_text, indent=None)[:2000]
    except Exception:
        refs_excerpt = "[]"

    sys_prompt += f"\nDesign Context (truncated): {design_excerpt}\nReference snippets: {refs_excerpt}\n"

    user_prompt = f"Create a PowerPoint plan for: {prompt}. Style/phase: {style}."
    if num_slides:
        user_prompt += f" Produce exactly {num_slides} slides."

    try:
        resp = text_client.chat.completions.create(
            model=CHAT_MODEL,
            messages=[{"role":"system","content":sys_prompt},{"role":"user","content":user_prompt}],
            max_tokens=1500,
            temperature=0.2
        )
    except Exception as e:
        logger.exception(f"Plan generation failed (API call): {e}")
        return [{"title":"Intro","bullets":["Overview"],"visual_required":False}]

    # DEBUG: print raw response for troubleshooting
    try:
        raw_text = resp.choices[0].message.content
        print("\n\n=== RAW GPT RESPONSE START ===\n")
        print(raw_text)
        print("\n=== RAW GPT RESPONSE END ===\n\n")
    except Exception as e:
        raw_text = ""
        print(f"\n\n=== ERROR READING RAW GPT RESPONSE: {e} ===\n\n")

    # try safe load first, then regex extraction
    plan = safe_json_load(raw_text)
    if not plan and raw_text:
        try:
            m = re.search(r'(\[.*\])', raw_text, re.DOTALL)
            if m:
                plan = json.loads(m.group(1))
        except Exception as e:
            logger.warning(f"Regex JSON extraction failed: {e}")

    if not plan:
        logger.warning(f"LLM returned invalid JSON (first 300 chars): {raw_text[:300]}")
        return [{"title":"Intro","bullets":["Overview"],"visual_required":False}]

    return plan

# ---------------------------------------------------------
# Image generation (image_client) - only when LLM requests visuals
# ---------------------------------------------------------
def generate_visual_image(visual_prompt: str):
    try:
        resp = image_client.images.generate(model=DALLE_MODEL, prompt=visual_prompt, size="1024x1024")
        # Support both url or b64 depending on response
        if getattr(resp, "data", None) and len(resp.data) > 0:
            # If URL provided
            item = resp.data[0]
            url = item.get("url") or item.get("b64_json")
            if not url:
                # fallback: try repr
                url = None
            # if b64_json, decode
            if item.get("b64_json"):
                import base64
                b = base64.b64decode(item["b64_json"])
                tmp = tempfile.NamedTemporaryFile(delete=False, suffix=".png")
                tmp.write(b)
                tmp.close()
                return tmp.name
            if url:
                r = requests.get(url, timeout=20)
                tmp = tempfile.NamedTemporaryFile(delete=False, suffix=".png")
                tmp.write(r.content)
                tmp.close()
                return tmp.name
        # fallback: no usable data
        return None
    except Exception as e:
        logger.warning(f"Visual generation failed for '{visual_prompt}': {e}")
        return None

# ---------------------------------------------------------
# PPT builder
# ---------------------------------------------------------
def build_ppt(slides):
    prs = Presentation()
    for s in slides:
        layout = prs.slide_layouts[1] if len(prs.slide_layouts) > 1 else prs.slide_layouts[0]
        slide = prs.slides.add_slide(layout)
        try:
            slide.shapes.title.text = s.get("title", "")
        except Exception:
            pass
        try:
            tf = slide.placeholders[1].text_frame
            tf.clear()
            for b in s.get("bullets", []):
                p = tf.add_paragraph()
                p.text = b
                p.font.size = Pt(18)
        except Exception:
            pass
        if s.get("image_path"):
            try:
                slide.shapes.add_picture(s["image_path"], Inches(0.5), Inches(3.0), width=Inches(8))
            except Exception as e:
                logger.debug(f"Failed to add visual to slide '{s.get('title')}': {e}")

    out = os.path.join(tempfile.gettempdir(), f"generated_presentation_{uuid.uuid4().hex[:8]}.pptx")
    prs.save(out)
    return out

# ---------------------------------------------------------
# Main pipeline: generate_presentation
# ---------------------------------------------------------
def generate_presentation(prompt: str,
                          style: str = "Auto",
                          requested_num_slides: int = None,
                          theme: str = None,
                          tag_filters: list = None):
    """
    Generate a new presentation.
    - prompt: user text
    - style: phase hint
    - requested_num_slides: override (optional)
    - theme: override (optional)
    - tag_filters: RAG tag filters (optional)
    """
    # auto-detect slide count & theme from prompt if not provided
    detected_slides, detected_theme = parse_user_intent(prompt)
    requested_num_slides = requested_num_slides or detected_slides
    theme = theme or detected_theme

    # 1) retrieve top-10 references via semantic search (text_client used internally)
    refs = semantic_search(prompt, top_k=10, tags=tag_filters) or []

    # 2) load design JSONs and short reference snippets
    design_context = []
    reference_texts = []
    for r in refs:
        ppt_name = r.get("ppt_name")
        text_snippet = (r.get("text") or "")[:400]
        if text_snippet:
            reference_texts.append(text_snippet)
        json_path = os.path.join("design_jsons", os.path.basename(ppt_name) + ".json")
        if os.path.exists(json_path):
            try:
                with open(json_path, "r", encoding="utf-8") as f:
                    design_context.append(json.load(f))
            except Exception:
                logger.debug(f"Failed to load design json for {ppt_name}")

    logger.info(f"Loaded {len(design_context)} design JSONs and {len(reference_texts)} reference snippets for context.")

    # 3) ask LLM for slide plan (may include visual decisions)
    plan = call_llm_plan(prompt, style, design_context, reference_texts,
                         num_slides=requested_num_slides, theme=theme)

    # If user requested exact slide count but LLM returned different count, try gentle adjustment
    if requested_num_slides and isinstance(plan, list) and len(plan) != requested_num_slides:
        try:
            adjust_prompt = (
                f"Adjust the following plan to exactly {requested_num_slides} slides. "
                "Keep intent and concise bullets; combine or split slides as needed.\n\n"
                f"Plan: {json.dumps(plan)[:4000]}"
            )
            adj_resp = text_client.chat.completions.create(
                model=CHAT_MODEL,
                messages=[{"role":"system","content":"You are a JSON-only assistant."},
                          {"role":"user","content":adjust_prompt}],
                max_tokens=800,
                temperature=0.2
            )
            adj_text = adj_resp.choices[0].message.content
            adj_plan = safe_json_load(adj_text)
            if not adj_plan and adj_text:
                m = re.search(r'(\[.*\])', adj_text, re.DOTALL)
                if m:
                    adj_plan = json.loads(m.group(1))
            if adj_plan and isinstance(adj_plan, list) and len(adj_plan) == requested_num_slides:
                plan = adj_plan
            else:
                logger.warning("Adjusted plan did not match requested length; using original plan.")
        except Exception:
            logger.exception("Failed to adjust plan to requested slide count; continuing with original plan.")

    # 4) build slides array, enrich with RAG bullets when appropriate, and generate visuals if LLM requested
    slides = []
    used_slide_ids = []

    # build a simple title->refs map for quick RAG mapping
    title_to_refs = {}
    for r in refs:
        t = (r.get("title") or "").strip()
        if not t:
            t = (r.get("text") or "").split("\n", 1)[0][:100]
        title_to_refs.setdefault(t.lower(), []).append(r)

    for p in plan:
        title = p.get("title", "Slide")
        bullets = p.get("bullets", []) or []
        refs_for_title = title_to_refs.get(title.lower(), refs)  # fallback to all refs
        rag_bullets = []
        for r in refs_for_title:
            txt = (r.get("text") or "").strip()
            if txt:
                rag_bullets.append(txt[:200])
            if len(rag_bullets) >= 3:
                break

        # prefer RAG bullets if available
        if rag_bullets:
            bullets = rag_bullets[:6]
        else:
            bullets = bullets[:6]

        slide = {"title": title, "bullets": bullets, "image_path": None}
        if p.get("visual_required"):
            visual_prompt = p.get("visual_prompt") or f"Professional illustration for slide titled '{title}'."
            img = generate_visual_image(visual_prompt)
            if img:
                slide["image_path"] = img

        slides.append(slide)

        # collect used slide ids for logging
        for r in (refs_for_title or []):
            sid = r.get("slide_id") or r.get("id")
            if sid:
                used_slide_ids.append(sid)

    # If requested_num_slides provided, enforce length (truncate or pad)
    if requested_num_slides:
        if len(slides) > requested_num_slides:
            slides = slides[:requested_num_slides]
        elif len(slides) < requested_num_slides:
            while len(slides) < requested_num_slides:
                slides.append({"title":"Additional","bullets":["(content to be filled)"], "image_path": None})

    # 5) build PPT, upload, and log
    out_path = build_ppt(slides)
    file_name = f"generated_{uuid.uuid4().hex[:8]}.pptx"
    try:
        upload_ppt_to_blob(out_path, file_name)
    except Exception:
        logger.exception("Failed to upload generated PPT to blob; continuing.")

    log = {
        "timestamp": now_ts(),
        "prompt": prompt,
        "style": style,
        "theme": theme,
        "requested_num_slides": requested_num_slides,
        "used_slide_ids": used_slide_ids,
        "retrieved_refs_count": len(refs),
        "retrieved_designs_count": len(design_context),
        "generated_file": file_name
    }
    try:
        upload_json_to_blob(json.dumps(log, indent=2).encode("utf-8"), blob_name=f"logs/{file_name}.json")
    except Exception:
        logger.exception("Failed to upload generation log to blob; continuing.")

    return out_path, log
