# ============================================================
# generate_ppt.py (FINAL UPDATED VERSION)
# ============================================================

import os
import tempfile
import uuid
import json
import re
import requests

from pptx import Presentation
from pptx.util import Inches, Pt

from utils import (
    get_env, safe_json_load, logger, now_ts,
    ensure_dir, text_client, image_client
)
from search_utils import semantic_search
from azure_blob_utils import upload_ppt_to_blob, upload_json_to_blob


ensure_dir("design_jsons")

CHAT_MODEL = get_env("CHAT_MODEL", required=True)
IMAGE_MODEL = get_env("IMAGE_MODEL", required=True)


# ---------------------------------------------------------
# Auto-detect slide count + theme
# ---------------------------------------------------------
def parse_user_intent(prompt: str):
    num_slides = None
    theme = None

    # detect number of slides ("5 slides")
    match = re.search(r"(\d+)\s+slides?", prompt.lower())
    if match:
        num_slides = int(match.group(1))

    # detect themes
    themes = ["corporate", "modern", "minimal", "professional",
              "dark", "light", "colorful", "gradient", "flat"]

    for t in themes:
        if t in prompt.lower():
            theme = t.capitalize()
            break

    return num_slides, theme


# ---------------------------------------------------------
# GPT PLAN GENERATOR
# ---------------------------------------------------------
def call_llm_plan(prompt, style, design_context, references_text,
                  num_slides=None, theme=None):

    sys_prompt = (
        "You are an expert presentation designer.\n"
        "Generate a slide plan in strict JSON format ONLY.\n"
        "Each item must be: {\n"
        "  \"title\": string,\n"
        "  \"bullets\": [strings],\n"
        "  \"visual_required\": boolean,\n"
        "  \"visual_prompt\": string\n"
        "}\n\n"
        "If the user explicitly requests images, visual_required MUST be true.\n"
        "Use reference text when helpful.\n"
        "Use design cues if available.\n"
    )

    if theme:
        sys_prompt += f"\nTheme preference: {theme}\n"

    sys_prompt += "\nDesign Context (truncated): " + json.dumps(design_context)[:2000]
    sys_prompt += "\nReference Text (truncated): " + json.dumps(references_text)[:1500]

    user_prompt = f"Create a PowerPoint plan for: {prompt}. Style: {style}."
    if num_slides:
        user_prompt += f" Produce exactly {num_slides} slides."

    try:
        resp = text_client.chat.completions.create(
            model=CHAT_MODEL,
            messages=[
                {"role": "system", "content": sys_prompt},
                {"role": "user", "content": user_prompt}
            ],
            max_completion_tokens=1500,
            temperature=1
        )

        raw = resp.choices[0].message.content
        plan = safe_json_load(raw)

        if not plan:
            logger.warning("Invalid plan JSON from GPT.")
            return [{"title": "Intro", "bullets": ["Overview"], "visual_required": False}]

        return plan

    except Exception as e:
        logger.exception(f"Plan generation failed: {e}")
        return [{"title": "Intro", "bullets": ["Overview"], "visual_required": False}]


# ---------------------------------------------------------
# IMAGE GENERATION
# ---------------------------------------------------------
def generate_visual_image(prompt: str):
    try:
        resp = image_client.images.generate(
            model=IMAGE_MODEL,
            prompt=prompt,
            size="1024x1024"
        )
        url = resp.data[0].url
        img_bytes = requests.get(url, timeout=20).content

        tmp = tempfile.NamedTemporaryFile(delete=False, suffix=".png")
        tmp.write(img_bytes)
        tmp.close()
        return tmp.name

    except Exception as e:
        logger.warning(f"Image generation failed: {e}")
        return None


# ---------------------------------------------------------
# PPT BUILDER  (UPDATED - FIXES LAYOUT)
# ---------------------------------------------------------
def build_ppt(slides):
    prs = Presentation()

    for s in slides:
        layout = prs.slide_layouts[1]  # Title + Body layout
        slide = prs.slides.add_slide(layout)

        # ---------------------------
        # TITLE
        # ---------------------------
        try:
            slide.shapes.title.text = s.get("title", "")
        except:
            pass

        # ---------------------------
        # IMAGE HANDLING
        # ---------------------------
        image_path = s.get("image_path")

        if image_path:
            # place the image on the right side
            img_left = Inches(6.5)
            img_top = Inches(1.5)
            img_width = Inches(2.5)

            try:
                slide.shapes.add_picture(image_path, img_left, img_top, width=img_width)
            except Exception as e:
                logger.debug(f"Failed to add image: {e}")
                image_path = None

        # ---------------------------
        # BODY/TEXT AREA
        # ---------------------------
        body = slide.placeholders[1].text_frame
        body.clear()

        textbox = slide.shapes[1]  # body placeholder shape reference

        if image_path:
            # Text uses left side only
            textbox.left = Inches(0.5)
            textbox.width = Inches(5.5)
        else:
            # Text expands full width
            textbox.left = Inches(0.5)
            textbox.width = Inches(9.0)

        # Add bullets
        for b in s.get("bullets", []):
            p = body.add_paragraph()
            p.text = b
            p.font.size = Pt(18)
            p.level = 0

    # Save file
    out_path = os.path.join(
        tempfile.gettempdir(),
        f"generated_presentation_{uuid.uuid4().hex[:8]}.pptx"
    )
    prs.save(out_path)
    return out_path


# ---------------------------------------------------------
# MAIN PIPELINE
# ---------------------------------------------------------
def generate_presentation(prompt: str, style="Auto", requested_num_slides=None,
                          theme=None, tag_filters=None, preview_only=False):

    detected_slides, detected_theme = parse_user_intent(prompt)

    requested_num_slides = requested_num_slides or detected_slides
    theme = theme or detected_theme

    # load references
    refs = semantic_search(prompt, top_k=5, tags=tag_filters) or []

    design_context = []
    reference_text = []

    for r in refs:
        ppt_name = r.get("ppt_name")
        text_snip = (r.get("text") or "")[:400]
        if text_snip:
            reference_text.append(text_snip)

        # load related design JSON if exists
        json_path = os.path.join("design_jsons", os.path.basename(ppt_name) + ".json")
        if os.path.exists(json_path):
            try:
                with open(json_path, "r", encoding="utf-8") as f:
                    design_context.append(json.load(f))
            except:
                pass

    logger.info(f"Loaded {len(design_context)} design JSONs and {len(reference_text)} text snippets")

    # ---------------------------
    # Generate PLAN
    # ---------------------------
    plan = call_llm_plan(
        prompt=prompt,
        style=style,
        design_context=design_context,
        references_text=reference_text,
        num_slides=requested_num_slides,
        theme=theme
    )

    # FAST PREVIEW MODE — return ONLY PLAN
    if preview_only:
        return None, {
            "preview_plan": plan,
            "slides": len(plan),
            "references_used": len(reference_text),
            "design_jsons_used": len(design_context),
        }

    # ---------------------------
    # Generate IMAGES + Assemble Slides
    # ---------------------------
    prepared_slides = []

    for sp in plan:
        slide = {
            "title": sp.get("title", "Untitled"),
            "bullets": sp.get("bullets", []),
            "image_path": None
        }

        if sp.get("visual_required"):
            img_prompt = sp.get("visual_prompt", f"Professional visual for {slide['title']}")
            slide["image_path"] = generate_visual_image(img_prompt)

        prepared_slides.append(slide)

    # Build PPT
    out_path = build_ppt(prepared_slides)

    # Upload PPT
    fname = f"generated_{uuid.uuid4().hex[:8]}.pptx"
    upload_ppt_to_blob(out_path, fname)

    # Log metadata
    log = {
        "timestamp": now_ts(),
        "prompt": prompt,
        "theme": theme,
        "requested_num_slides": requested_num_slides,
        "references_used": len(reference_text),
        "design_jsons_used": len(design_context),
        "slides_generated": len(prepared_slides),
        "ppt_file": fname
    }

    upload_json_to_blob(json.dumps(log, indent=2).encode("utf-8"),
                        f"logs/{fname}.json")

    return out_path, log
