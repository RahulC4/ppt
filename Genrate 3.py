# generate_ppt.py (robust version – forces visuals, proper slide count, stronger plan generation)
import os
import tempfile
import uuid
import json
import re
import requests
from pptx import Presentation
from pptx.util import Inches, Pt
from openai import AzureOpenAI
from utils import get_env, safe_json_load, logger, now_ts
from search_utils import semantic_search
from azure_blob_utils import upload_ppt_to_blob, upload_json_to_blob


# --- Azure OpenAI Clients ---
text_client = AzureOpenAI(
    azure_endpoint=get_env("OPENAI_API_BASE", required=True),
    api_key=get_env("OPENAI_API_KEY", required=True),
    api_version=get_env("OPENAI_API_VERSION", "2024-05-01-preview")
)

image_client = AzureOpenAI(
    azure_endpoint=get_env("IMAGE_API_BASE", required=True),
    api_key=get_env("IMAGE_API_KEY", required=True),
    api_version=get_env("OPENAI_API_VERSION", "2024-05-01-preview")
)

CHAT_MODEL = get_env("CHAT_MODEL", "gpt-5-mini")
DALLE_MODEL = get_env("IMAGE_MODEL", "gpt-image-1-mini")


# ---------------------------------------------------------
# Extract number of slides + theme
# ---------------------------------------------------------
def parse_user_intent(prompt: str):
    num_slides = None
    theme = None

    # detect numbers like "5 slides"
    match = re.search(r'(\d+)\s*slides?', prompt.lower())
    if match:
        num_slides = int(match.group(1))

    # detect theme keywords
    for word in ["modern", "minimal", "corporate", "professional", "dark", "light", "colorful"]:
        if word in prompt.lower():
            theme = word.capitalize()
            break

    return num_slides, theme


# ---------------------------------------------------------
# GPT Plan generator — upgraded prompt
# ---------------------------------------------------------
def call_llm_plan(prompt, design_context, references_text, requested_slides=None, theme=None):
    sys_prompt = (
        "You are an expert presentation designer.\n\n"
        "You MUST output STRICT JSON: a list of slides.\n"
        "Each slide MUST contain these keys:\n"
        "{\n"
        "  \"title\": string,\n"
        "  \"bullets\": list of at least 3 strings,\n"
        "  \"visual_required\": boolean,\n"
        "  \"visual_prompt\": string (REQUIRED if visual_required=true)\n"
        "}\n\n"
        "Rules:\n"
        "1. ALWAYS produce the exact slide count requested.\n"
        "2. If not specified, choose a professional slide count (5–8).\n"
        "3. Use design JSONs to match style, layouts, shapes, accents.\n"
        "4. If any reference slide uses images/icons/shapes, set visual_required=true.\n"
        "5. NEVER output fewer than 3 bullets per slide.\n"
        "6. NEVER output commentary—ONLY pure JSON.\n"
    )

    if theme:
        sys_prompt += f"\nTheme: {theme}.\n"

    design_excerpt = json.dumps(design_context, ensure_ascii=False)[:3000]
    refs_excerpt = json.dumps(references_text, ensure_ascii=False)[:1500]

    sys_prompt += f"\nDesign Context (truncated): {design_excerpt}\n"
    sys_prompt += f"\nReference Text: {refs_excerpt}\n"

    user_prompt = f"Create a slide plan for: {prompt}."
    if requested_slides:
        user_prompt += f" Produce exactly {requested_slides} slides."

    try:
        resp = text_client.chat.completions.create(
            model=CHAT_MODEL,
            messages=[
                {"role": "system", "content": sys_prompt},
                {"role": "user", "content": user_prompt}
            ],
            max_tokens=1000
        )
        text = resp.choices[0].message.content
        plan = safe_json_load(text)
        return plan
    except Exception as e:
        logger.exception(f"LLM plan generation failed: {e}")
        return None


# ---------------------------------------------------------
# Visual generation fallback
# ---------------------------------------------------------
def generate_visual_image(visual_prompt):
    try:
        resp = image_client.images.generate(
            model=DALLE_MODEL,
            prompt=visual_prompt,
            size="1024x1024"
        )
        url = resp.data[0].url
        img_data = requests.get(url, timeout=20).content
        tmp = tempfile.NamedTemporaryFile(delete=False, suffix=".png")
        tmp.write(img_data)
        tmp.close()
        return tmp.name
    except Exception as e:
        logger.warning(f"Visual generation failed: {e}")
        return None


# ---------------------------------------------------------
# PPT builder
# ---------------------------------------------------------
def build_ppt(slides):
    prs = Presentation()
    for s in slides:
        layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(layout)

        # title
        try:
            slide.shapes.title.text = s.get("title", "")
        except:
            pass

        # bullets
        try:
            tf = slide.placeholders[1].text_frame
            tf.clear()
            for b in s.get("bullets", []):
                p = tf.add_paragraph()
                p.text = b
                p.font.size = Pt(18)
        except:
            pass

        # image
        if s.get("image_path"):
            try:
                slide.shapes.add_picture(s["image_path"], Inches(0.5), Inches(3), width=Inches(8))
            except:
                pass

    out = os.path.join(tempfile.gettempdir(), f"ppt_{uuid.uuid4().hex[:8]}.pptx")
    prs.save(out)
    return out


# ---------------------------------------------------------
# MAIN PIPELINE
# ---------------------------------------------------------
def generate_presentation(prompt: str, tag_filters=None):
    detected_slides, detected_theme = parse_user_intent(prompt)

    # semantic search
    refs = semantic_search(prompt, top_k=10, tags=tag_filters) or []

    # gather design + text
    design_context, ref_texts = [], []
    for r in refs:
        ppt_name = r.get("ppt_name")
        text = (r.get("text") or "")[:400]
        if text:
            ref_texts.append(text)

        json_path = os.path.join("design_jsons", f"{ppt_name}.json")
        if os.path.exists(json_path):
            try:
                with open(json_path, "r", encoding="utf-8") as f:
                    design_context.append(json.load(f))
            except:
                pass

    logger.info(f"Loaded {len(design_context)} design JSONs and {len(ref_texts)} reference snippets.")

    # slide plan
    plan = call_llm_plan(
        prompt,
        design_context=design_context,
        references_text=ref_texts,
        requested_slides=detected_slides,
        theme=detected_theme
    )

    if not plan:
        plan = []

    # force correct slide count
    final_slides = []
    required = detected_slides or 5
    for s in plan[:required]:
        s["visual_required"] = s.get("visual_required", True)
        if s["visual_required"] and not s.get("visual_prompt"):
            s["visual_prompt"] = f"Professional illustration for {s.get('title', '')}"
        final_slides.append(s)

    while len(final_slides) < required:
        final_slides.append({
            "title": f"Slide {len(final_slides)+1}",
            "bullets": ["Point 1", "Point 2", "Point 3"],
            "visual_required": True,
            "visual_prompt": "Professional corporate slide visual"
        })

    # generate images
    for s in final_slides:
        if s.get("visual_required"):
            img = generate_visual_image(s["visual_prompt"])
            s["image_path"] = img

    # build ppt
    ppt_path = build_ppt(final_slides)
    file_name = f"generated_{uuid.uuid4().hex[:8]}.pptx"
    upload_ppt_to_blob(ppt_path, file_name)

    # log upload
    log = {
        "prompt": prompt,
        "slides": len(final_slides),
        "design_jsons_used": len(design_context),
        "refs_used": len(ref_texts),
        "file": file_name
    }
    upload_json_to_blob(json.dumps(log).encode(), f"logs/{file_name}.json")

    return ppt_path, log
